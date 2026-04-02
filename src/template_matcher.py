"""
슬라이드 ↔ 템플릿 매칭 엔진.
새 슬라이드의 shape 구성을 분석하여 기존 T0~T9 템플릿과 유사도를 비교.
"""
import json
from collections import Counter
from pptx import Presentation


TEMPLATE_NAMES = {
    'T0': '구분페이지',
    'T1': '카드형 다중 (현황분석)',
    'T2': '카드+다이어그램 (목적/전략)',
    'T3': '범위/개요 (거버닝메시지)',
    'T4': '다중 데이터테이블',
    'T5': '테이블+다이어그램',
    'T6': '순수 데이터테이블',
    'T7': '테이블+설명shape (프로세스)',
    'T8': '이미지중심',
    'T9': '핵심메시지/다이어그램',
    'T14': '기타/특수',
}


def extract_slide_features(slide):
    """슬라이드의 shape 구성 특성을 추출"""
    auto_shapes = 0
    text_boxes = 0
    tables = 0
    card_tables = 0  # 2x1 테이블
    data_tables = 0
    pictures = 0
    groups = 0
    total_text_len = 0

    for shape in slide.shapes:
        st = str(shape.shape_type)
        if 'AUTO' in st:
            auto_shapes += 1
        elif 'TEXT' in st:
            text_boxes += 1
        elif 'TABLE' in st:
            tables += 1
            t = shape.table
            nr = len(list(t.rows))
            nc = len(t.columns)
            if nr == 2 and nc == 1:
                card_tables += 1
            else:
                data_tables += 1
        elif 'PICTURE' in st:
            pictures += 1
        elif 'GROUP' in st:
            groups += 1

        if shape.has_text_frame:
            total_text_len += len(shape.text_frame.text.strip())

    total = len(slide.shapes)

    return {
        'total_shapes': total,
        'auto_shapes': auto_shapes,
        'text_boxes': text_boxes,
        'tables': tables,
        'card_tables': card_tables,
        'data_tables': data_tables,
        'pictures': pictures,
        'groups': groups,
        'total_text_len': total_text_len,
    }


# 각 템플릿 타입의 대표 특성 (평균)
TEMPLATE_PROFILES = {
    'T0': {'total_shapes': (3, 8), 'card_tables': 0, 'data_tables': (0, 1), 'pictures': 0, 'auto_shapes': (0, 2)},
    'T1': {'total_shapes': (9, 50), 'card_tables': (2, 6), 'data_tables': 0, 'pictures': (0, 6)},
    'T2': {'total_shapes': (10, 20), 'card_tables': (2, 3), 'data_tables': 0, 'auto_shapes': (3, 10)},
    'T3': {'total_shapes': (10, 20), 'card_tables': (0, 2), 'data_tables': 0, 'text_boxes': (3, 8)},
    'T4': {'total_shapes': (10, 50), 'card_tables': 0, 'data_tables': (2, 5)},
    'T5': {'total_shapes': (15, 70), 'card_tables': 0, 'data_tables': 1, 'auto_shapes': (10, 30)},
    'T6': {'total_shapes': (3, 20), 'card_tables': 0, 'data_tables': 1, 'auto_shapes': (0, 3)},
    'T7': {'total_shapes': (10, 40), 'card_tables': 0, 'data_tables': 1, 'auto_shapes': (3, 20)},
    'T8': {'total_shapes': (10, 130), 'pictures': (2, 20)},
    'T9': {'total_shapes': (14, 95), 'card_tables': 0, 'data_tables': 0, 'auto_shapes': (5, 30)},
}


def _in_range(value, spec):
    """값이 스펙 범위 내인지 확인. spec이 int면 정확 매치, tuple이면 범위"""
    if isinstance(spec, tuple):
        return spec[0] <= value <= spec[1]
    return value == spec


def match_template(features: dict) -> list:
    """
    슬라이드 특성을 기존 템플릿과 비교하여 유사도 순으로 반환.
    Returns: [{'template': 'T1', 'score': 0.85, 'name': '카드형 다중'}, ...]
    """
    results = []

    for tmpl, profile in TEMPLATE_PROFILES.items():
        score = 0
        total_checks = 0

        for key, spec in profile.items():
            if key in features:
                total_checks += 1
                if _in_range(features[key], spec):
                    score += 1

        if total_checks > 0:
            similarity = score / total_checks
        else:
            similarity = 0

        results.append({
            'template': tmpl,
            'score': round(similarity, 2),
            'name': TEMPLATE_NAMES.get(tmpl, '기타'),
        })

    results.sort(key=lambda x: x['score'], reverse=True)
    return results


def analyze_and_match(pptx_path: str, slide_number: int) -> dict:
    """PPTX 파일의 특정 슬라이드를 분석하여 템플릿 매칭"""
    prs = Presentation(pptx_path)

    if slide_number < 1 or slide_number > len(prs.slides):
        return {'error': f'슬라이드 번호 {slide_number}이 범위를 벗어났습니다 (1~{len(prs.slides)})'}

    slide = prs.slides[slide_number - 1]
    features = extract_slide_features(slide)
    matches = match_template(features)

    return {
        'slide_number': slide_number,
        'features': features,
        'matches': matches[:5],  # 상위 5개
        'best_match': matches[0] if matches else None,
    }
