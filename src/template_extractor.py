"""
참조 PPTX를 분석하여 slide_index.json을 생성하는 도구.
각 슬라이드의 shape별 역할(제목, 본문, 카드, 테이블 등)을 자동 분류.
"""
import json
import sys
import os
from pptx import Presentation
from pptx.util import Emu


# 템플릿 분류 맵 (분석 결과 기반)
TEMPLATE_MAP = {
    1:'T3', 2:'T3', 3:'T1', 4:'T0', 5:'T1', 6:'T1', 7:'T1', 8:'T1',
    9:'T0', 10:'T0', 11:'T14', 12:'T14', 13:'T8', 14:'T14', 15:'T4',
    16:'T9', 17:'T0', 18:'T1', 19:'T0', 20:'T7', 21:'T0', 22:'T6',
    23:'T6', 24:'T7', 25:'T14', 26:'T4', 27:'T8', 28:'T7', 29:'T9',
    30:'T7', 31:'T9', 32:'T0', 33:'T7', 34:'T0', 35:'T0', 36:'T0',
    37:'T0', 38:'T0', 39:'T0', 40:'T8', 41:'T14', 42:'T0', 43:'T5',
    44:'T5', 45:'T4', 46:'T4', 47:'T9', 48:'T6', 49:'T4', 50:'T4',
    51:'T6', 52:'T4', 53:'T5', 54:'T6', 55:'T8', 56:'T0', 57:'T6',
    58:'T14', 59:'T7', 60:'T0', 61:'T6', 62:'T9', 63:'T14', 64:'T7',
    65:'T0', 66:'T6', 67:'T6', 68:'T6', 69:'T6', 70:'T6', 71:'T9',
    72:'T9', 73:'T4', 74:'T4', 75:'T0', 76:'T8', 77:'T2', 78:'T1',
    79:'T6', 80:'T7', 81:'T4', 82:'T9', 83:'T14', 84:'T4', 85:'T0',
    86:'T14', 87:'T7', 88:'T0', 89:'T5', 90:'T7', 91:'T8', 92:'T5',
    93:'T8', 94:'T0', 95:'T0', 96:'T5', 97:'T7', 98:'T8', 99:'T4',
    100:'T7', 101:'T14', 102:'T5', 103:'T14', 104:'T4', 105:'T7',
    106:'T4', 107:'T5', 108:'T9', 109:'T4', 110:'T6', 111:'T9',
    112:'T4', 113:'T4', 114:'T14', 115:'T7', 116:'T8', 117:'T5',
    118:'T6', 119:'T9', 120:'T14', 121:'T9', 122:'T5', 123:'T4',
    124:'T6', 125:'T8', 126:'T5', 127:'T4', 128:'T4'
}

# III권 템플릿 분류 맵 (47 슬라이드)
TEMPLATE_MAP_VOL3 = {
    1:'T9', 2:'T9', 3:'T0', 4:'T0', 5:'T0', 6:'T1', 7:'T0', 8:'T0',
    9:'T0', 10:'T9', 11:'T0', 12:'T5', 13:'T0', 14:'T0', 15:'T1',
    16:'T0', 17:'T0', 18:'T6', 19:'T6', 20:'T0', 21:'T0', 22:'T8',
    23:'T0', 24:'T6', 25:'T8', 26:'T0', 27:'T9', 28:'T6', 29:'T0',
    30:'T3', 31:'T0', 32:'T0', 33:'T8', 34:'T6', 35:'T0', 36:'T0',
    37:'T6', 38:'T0', 39:'T0', 40:'T8', 41:'T7', 42:'T0', 43:'T0',
    44:'T0', 45:'T0', 46:'T9', 47:'T9'
}

TEMPLATE_NAMES = {
    'T0': '구분페이지',
    'T1': '카드형 다중 (현황분석)',
    'T2': '카드+다이어그램 (목적/전략)',
    'T3': '범위/개요 (거버닝메시지)',
    'T4': '다중 데이터테이블',
    'T5': '테이블+다이어그램',
    'T6': '순수 데이터테이블',
    'T7': '테이블+설명shape (프로세스)',
    'T8': '이미지중심',
    'T9': '핵심메시지/다이어그램',
    'T14': '기타/특수',
}


def classify_shape_role(shape, slide_shapes_context):
    """shape의 역할을 자동 분류"""
    name = shape.name
    shape_type = str(shape.shape_type)

    # 테이블
    if shape.has_table:
        table = shape.table
        nrows = len(list(table.rows))
        ncols = len(table.columns)
        if nrows == 2 and ncols == 1:
            return 'card_table'  # 2x1 카드형 테이블
        return 'data_table'

    # 이미지
    if 'PICTURE' in shape_type:
        return 'image'

    # 그룹 (보통 배경 장식)
    if 'GROUP' in shape_type:
        return 'group_decoration'

    # 라인/프리폼
    if 'LINE' in shape_type or 'FREE' in shape_type:
        return 'decoration'

    # 텍스트가 있는 shape
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if not text:
            return 'empty_shape'

        # Governing Message
        if 'Governing' in text and 'Message' in text:
            return 'governing_label'

        # 제목 placeholder
        if name.startswith('제목'):
            return 'breadcrumb'  # 섹션 경로

        # 부제목 placeholder
        if name.startswith('부제목'):
            return 'governing_message'  # 실제 거버닝 메시지 내용

        # Chapter 표시
        if text.strip() == 'Chapter':
            return 'chapter_label'

        # 페이지 번호
        if '‹#›' in text or text.startswith('Ⅱ -') or text.startswith('Ⅲ -'):
            return 'page_number'

        # 섹션 타이틀 (■ 마커)
        if name.startswith('TextBox') and len(text) < 50:
            return 'section_title'

        # 둥근 사각형 계열
        if '모서리가 둥근' in name or '둥근 모서리' in name:
            if len(text) < 30:
                return 'label_box'  # 라벨
            return 'content_box'  # 내용 박스

        # 직사각형 계열
        if name.startswith('직사각형'):
            if len(text) < 40:
                return 'heading_box'  # 소제목 박스
            return 'content_box'

        # 양쪽 모서리가 둥근 사각형
        if '양쪽 모서리' in name:
            return 'content_box'

        # AutoShape 계열
        if 'AUTO' in shape_type or 'Google Shape' in name or 'AutoShape' in name:
            if len(text) < 40:
                return 'label_shape'
            return 'content_shape'

        # Rectangle 계열
        if name.startswith('Rectangle') or name.startswith('사각형'):
            if len(text) < 40:
                return 'heading_box'
            return 'content_box'

        # 타원
        if name.startswith('타원'):
            return 'number_circle'

        # Text Box
        if 'Text Box' in name or 'TextBox' in name:
            if len(text) < 30:
                return 'section_title'
            return 'text_content'

        return 'text_content'

    return 'unknown'


def extract_shape_info(shape):
    """shape의 기본 정보 추출"""
    info = {
        'name': shape.name,
        'shape_type': str(shape.shape_type),
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
    }

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        info['text'] = text[:200] if text else ''
        info['text_length'] = len(text)

    if shape.has_table:
        table = shape.table
        nrows = len(list(table.rows))
        ncols = len(table.columns)
        info['table_size'] = f'{nrows}x{ncols}'
        # 첫 몇 셀 내용
        cells_preview = []
        for ri, row in enumerate(table.rows):
            if ri >= 3:
                break
            for ci, cell in enumerate(row.cells):
                if ci >= 3:
                    break
                cells_preview.append(cell.text.strip()[:50])
        info['table_preview'] = cells_preview

    return info


def extract_slide_index(pptx_path, output_path, template_map=None,
                        slide_offset=0, source_label=""):
    """참조 PPTX를 분석하여 slide_index.json 생성.

    Args:
        template_map: 사용할 템플릿 맵 (기본: TEMPLATE_MAP)
        slide_offset: slide_number에 더할 오프셋 (merge용)
        source_label: source_pptx 필드에 저장할 라벨 (빈 문자열이면 파일명 사용)
    """
    if template_map is None:
        template_map = TEMPLATE_MAP

    prs = Presentation(pptx_path)
    source_name = source_label or os.path.basename(pptx_path)

    slide_index = {
        'source_pptx': source_name,
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height,
        'total_slides': len(prs.slides),
        'template_types': TEMPLATE_NAMES,
        'slides': []
    }

    for idx, slide in enumerate(prs.slides):
        actual_num = idx + 1
        slide_num = actual_num + slide_offset
        template_code = template_map.get(actual_num, 'T14')

        shapes_info = []
        for shape in slide.shapes:
            role = classify_shape_role(shape, slide.shapes)
            s_info = extract_shape_info(shape)
            s_info['role'] = role
            shapes_info.append(s_info)

        # 역할별 shape 인덱스 매핑 (텍스트 교체 시 사용)
        role_map = {}
        for si, s in enumerate(shapes_info):
            role = s['role']
            if role not in role_map:
                role_map[role] = []
            role_map[role].append(si)

        slide_entry = {
            'slide_number': slide_num,
            'source_pptx': source_name,
            'source_slide': actual_num,
            'template': template_code,
            'template_name': TEMPLATE_NAMES.get(template_code, '기타'),
            'layout_name': slide.slide_layout.name,
            'shape_count': len(shapes_info),
            'role_map': role_map,
            'shapes': shapes_info,
            'image_file': f'S{slide_num:04d}_{template_code}.png'
        }

        slide_index['slides'].append(slide_entry)

        if actual_num % 20 == 0:
            print(f'  Processed {actual_num}/{len(prs.slides)}...')

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(slide_index, f, ensure_ascii=False, indent=2)

    print(f'Saved: {output_path}')
    print(f'Total slides: {len(slide_index["slides"])}')

    # 템플릿별 통계
    from collections import Counter
    tmpl_counts = Counter(s['template'] for s in slide_index['slides'])
    print('\nTemplate distribution:')
    for tmpl, count in tmpl_counts.most_common():
        print(f'  {tmpl} ({TEMPLATE_NAMES.get(tmpl, "?")}): {count}')

    return slide_index


def extract_and_merge(pptx_list, output_path):
    """여러 PPTX를 분석하여 하나의 slide_index.json으로 병합.

    Args:
        pptx_list: [(pptx_path, template_map, slide_offset, source_label), ...]
        output_path: 병합된 slide_index.json 저장 경로
    """
    merged_slides = []
    first_width = None
    first_height = None

    for pptx_path, template_map, slide_offset, source_label in pptx_list:
        print(f'\n=== Analyzing: {os.path.basename(pptx_path)} (offset={slide_offset}) ===')
        prs = Presentation(pptx_path)
        source_name = source_label or os.path.basename(pptx_path)

        if first_width is None:
            first_width = prs.slide_width
            first_height = prs.slide_height

        for idx, slide in enumerate(prs.slides):
            actual_num = idx + 1
            slide_num = actual_num + slide_offset
            template_code = template_map.get(actual_num, 'T14')

            shapes_info = []
            for shape in slide.shapes:
                role = classify_shape_role(shape, slide.shapes)
                s_info = extract_shape_info(shape)
                s_info['role'] = role
                shapes_info.append(s_info)

            role_map = {}
            for si, s in enumerate(shapes_info):
                role = s['role']
                if role not in role_map:
                    role_map[role] = []
                role_map[role].append(si)

            slide_entry = {
                'slide_number': slide_num,
                'source_pptx': source_name,
                'source_slide': actual_num,
                'template': template_code,
                'template_name': TEMPLATE_NAMES.get(template_code, '기타'),
                'layout_name': slide.slide_layout.name,
                'shape_count': len(shapes_info),
                'role_map': role_map,
                'shapes': shapes_info,
                'image_file': f'S{slide_num:04d}_{template_code}.png'
            }
            merged_slides.append(slide_entry)

            if actual_num % 20 == 0:
                print(f'  Processed {actual_num}/{len(prs.slides)}...')

        print(f'  Done: {len(prs.slides)} slides from {source_name}')

    merged_index = {
        'source_pptx': 'merged',
        'slide_width': first_width,
        'slide_height': first_height,
        'total_slides': len(merged_slides),
        'template_types': TEMPLATE_NAMES,
        'slides': merged_slides
    }

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(merged_index, f, ensure_ascii=False, indent=2)

    print(f'\nMerged saved: {output_path}')
    print(f'Total slides: {len(merged_slides)}')

    from collections import Counter
    tmpl_counts = Counter(s['template'] for s in merged_slides)
    print('\nTemplate distribution:')
    for tmpl, count in tmpl_counts.most_common():
        print(f'  {tmpl} ({TEMPLATE_NAMES.get(tmpl, "?")}): {count}')

    return merged_index


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage:')
        print('  python template_extractor.py <pptx_path> [output_json_path]')
        print('  python template_extractor.py merge <vol2.pptx> <vol3.pptx> [output_json_path]')
        sys.exit(1)

    if sys.argv[1] == 'merge':
        if len(sys.argv) < 4:
            print('Usage: python template_extractor.py merge <vol2.pptx> <vol3.pptx> [output_json_path]')
            sys.exit(1)
        vol2_path = sys.argv[2]
        vol3_path = sys.argv[3]
        output_path = sys.argv[4] if len(sys.argv) > 4 else 'templates/slide_index.json'

        pptx_list = [
            (vol2_path, TEMPLATE_MAP, 2000, os.path.basename(vol2_path)),
            (vol3_path, TEMPLATE_MAP_VOL3, 3000, os.path.basename(vol3_path)),
        ]
        extract_and_merge(pptx_list, output_path)
    else:
        pptx_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else 'templates/slide_index.json'
        extract_slide_index(pptx_path, output_path)
