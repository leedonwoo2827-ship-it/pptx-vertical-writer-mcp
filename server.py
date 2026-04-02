# -*- coding: utf-8 -*-
"""
md2ppt MCP Server

참조 PPTX 템플릿 기반으로 마크다운을 PowerPoint 프레젠테이션으로 변환하는 MCP 서버.
회사 표준 템플릿의 레이아웃/디자인을 100% 보존하면서 텍스트만 교체합니다.

실행: python server.py
"""

import json
import sys
from pathlib import Path

BASE_DIR = Path(__file__).parent
sys.path.insert(0, str(BASE_DIR / "src"))

from mcp.server.fastmcp import FastMCP
from md_parser import parse_md
from slide_builder import build_presentation, load_slide_index

mcp = FastMCP("pptx-vertical-writer")

DEFAULT_TEMPLATE_DIR = BASE_DIR / "templates"
DEFAULT_SLIDE_INDEX = DEFAULT_TEMPLATE_DIR / "slide_index.json"
DEFAULT_PLACEHOLDER_PPTX = DEFAULT_TEMPLATE_DIR / "placeholder.pptx"
TEMPLATE_IMAGES_DIR = BASE_DIR / "template-images" / "all-slides"


# ---------------------------------------------------------------------------
# 내부 유틸리티
# ---------------------------------------------------------------------------

def _resolve_output(project_dir: str, output_file: str, fallback: str = "output.pptx") -> Path:
    if project_dir:
        proj = Path(project_dir)
        out_dir = proj / "output"
        out_dir.mkdir(parents=True, exist_ok=True)
        name = Path(output_file).name if output_file else fallback
        return out_dir / name
    if output_file:
        p = Path(output_file)
        return p if p.is_absolute() else Path.home() / "Documents" / output_file
    return Path.home() / "Documents" / fallback


# ---------------------------------------------------------------------------
# MCP 도구
# ---------------------------------------------------------------------------

@mcp.tool()
def list_templates() -> str:
    """사용 가능한 슬라이드 템플릿 목록을 반환합니다.

    각 템플릿 타입(T0~T9)의 용도, 슬라이드 수, 대표 이미지 경로를 포함합니다.
    글쓰기 스킬에서 적합한 템플릿을 추천할 때 사용합니다.

    Returns:
        템플릿 카탈로그 (JSON)
    """
    if not DEFAULT_SLIDE_INDEX.exists():
        return "오류: slide_index.json이 없습니다. analyze_template()을 먼저 실행하세요."

    with open(DEFAULT_SLIDE_INDEX, 'r', encoding='utf-8') as f:
        idx = json.load(f)

    # 템플릿별 통계
    from collections import Counter
    tmpl_counts = Counter(s['template'] for s in idx.get('slides', []))

    TEMPLATE_NAMES = {
        'T0': '구분페이지 — 섹션 구분, 챕터 시작. 글: 제목 1줄',
        'T1': '카드형 다중 — 현황/문제점/개선, 전략/방법론. 글: 카드 2~6개(제목+내용)',
        'T2': '카드+다이어그램 — 목적/전략 카드 + 세분화. 글: 카드 + 불릿 리스트',
        'T3': '범위/개요 — Governing Message + 둥근사각형 그리드. 글: 핵심 문구 + 6개 영역',
        'T4': '다중 데이터테이블 — 복수 테이블 + 일정표. 글: 마크다운 테이블',
        'T5': '테이블+다이어그램 — 단일 테이블 + shape 설명. 글: 테이블 + 설명',
        'T6': '순수 데이터테이블 — 큰 테이블 중심. 글: 마크다운 테이블',
        'T7': '테이블+설명shape — 프로세스 화살표 + 테이블. 글: 단계별 설명 + 테이블',
        'T8': '이미지중심 — 조직도, 스크린샷, 구성도. 글: 이미지 경로 + 캡션',
        'T9': '핵심메시지/다이어그램 — 둥근사각형 라벨, 프로세스 흐름도. 글: 핵심 문구 리스트',
        'T14': '기타/특수 — 분류 외 레이아웃',
    }

    catalog = []
    for tmpl_code in sorted(tmpl_counts.keys()):
        # 대표 슬라이드 번호
        rep_slides = [s['slide_number'] for s in idx['slides'] if s['template'] == tmpl_code][:3]
        # 대표 이미지 경로
        images = []
        for sn in rep_slides:
            img = TEMPLATE_IMAGES_DIR / f"S{sn:03d}_{tmpl_code}.png"
            if img.exists():
                images.append(str(img))

        catalog.append({
            'template': tmpl_code,
            'name': TEMPLATE_NAMES.get(tmpl_code, '기타'),
            'slide_count': tmpl_counts[tmpl_code],
            'representative_slides': rep_slides,
            'preview_images': images,
        })

    return json.dumps(catalog, ensure_ascii=False, indent=2)


@mcp.tool()
def analyze_template(
    pptx_path: str,
    output_dir: str = ""
) -> str:
    """참조 PPTX를 분석하여 slide_index.json과 placeholder 템플릿을 생성합니다.

    새로운 회사 템플릿을 등록할 때 1회 실행합니다.
    참조 PPTX의 각 슬라이드를 분석하여 shape별 역할(제목/본문/카드/테이블)을 자동 분류하고,
    텍스트를 플레이스홀더(████/Lorem)로 치환한 깨끗한 템플릿을 생성합니다.

    Args:
        pptx_path: 분석할 참조 PPTX 파일 경로
        output_dir: 출력 디렉토리 (생략 시 templates/)

    Returns:
        분석 결과 요약
    """
    from template_extractor import extract_slide_index
    from template_sanitizer import sanitize_pptx, sanitize_slide_index

    pptx = Path(pptx_path)
    if not pptx.exists():
        return f"오류: PPTX 파일을 찾을 수 없습니다: {pptx_path}"

    out = Path(output_dir) if output_dir else DEFAULT_TEMPLATE_DIR
    out.mkdir(parents=True, exist_ok=True)

    idx_path = str(out / "slide_index.json")
    placeholder_path = str(out / "placeholder.pptx")

    try:
        # 1. 슬라이드 분석 → slide_index.json
        slide_index = extract_slide_index(str(pptx), idx_path)

        # 2. 플레이스홀더 PPTX 생성
        sanitize_pptx(str(pptx), idx_path, placeholder_path)

        # 3. slide_index.json 텍스트도 치환
        sanitize_slide_index(idx_path)

        total = slide_index['total_slides']
        from collections import Counter
        dist = Counter(s['template'] for s in slide_index['slides'])

        result = f"분석 완료!\n"
        result += f"  슬라이드: {total}장\n"
        result += f"  slide_index: {idx_path}\n"
        result += f"  placeholder: {placeholder_path}\n"
        result += f"\n템플릿 분포:\n"
        for t, c in dist.most_common():
            result += f"  {t}: {c}장\n"

        return result

    except Exception as e:
        return f"오류: {type(e).__name__}: {e}"


@mcp.tool()
def create_pptx(
    extended_md: str,
    output_file: str = "",
    project_dir: str = "",
    template_dir: str = ""
) -> str:
    """확장 마크다운으로부터 PowerPoint 프레젠테이션을 생성합니다.

    확장 MD 포맷:
    ---config
    reference_pptx: path/to/placeholder.pptx  (생략 시 기본 템플릿 사용)
    ---

    ---slide
    template: T1
    ref_slide: 5
    ---
    @governing_message: 거버닝 메시지 텍스트
    @카드1_제목: 첫 번째 카드 제목
    @카드1_내용: 첫 번째 카드 내용
    @content_1: 본문 텍스트

    Args:
        extended_md: 확장 마크다운 텍스트 (---slide 구분자 포함)
        output_file: 출력 PPTX 파일명
        project_dir: 프로젝트 폴더 (지정 시 output/ 하위에 저장)
        template_dir: 템플릿 디렉토리 (생략 시 기본 templates/)

    Returns:
        생성된 파일 경로
    """
    tmpl_dir = Path(template_dir) if template_dir else DEFAULT_TEMPLATE_DIR
    idx_path = tmpl_dir / "slide_index.json"
    placeholder = tmpl_dir / "placeholder.pptx"

    if not idx_path.exists():
        return "오류: slide_index.json이 없습니다. analyze_template()을 먼저 실행하세요."

    out_path = _resolve_output(project_dir, output_file)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        # MD 파싱
        md_data = parse_md(extended_md)

        # reference_pptx 기본값 설정
        if 'reference_pptx' not in md_data.get('config', {}):
            md_data.setdefault('config', {})['reference_pptx'] = str(placeholder)

        config_ref = md_data['config'].get('reference_pptx', '')
        if not Path(config_ref).is_absolute():
            md_data['config']['reference_pptx'] = str(BASE_DIR / config_ref)

        slide_index = load_slide_index(str(idx_path))

        # PPTX 생성
        build_presentation(md_data, slide_index, str(out_path))

        if out_path.exists():
            size = out_path.stat().st_size
            slide_count = len(md_data.get('slides', []))
            return f"PPTX 생성 완료!\n파일: {out_path}\n크기: {size:,} bytes\n슬라이드: {slide_count}장"
        return f"오류: 파일 생성 실패: {out_path}"

    except Exception as e:
        return f"오류: {type(e).__name__}: {e}"


# ---------------------------------------------------------------------------
# 템플릿 관리 도구
# ---------------------------------------------------------------------------

@mcp.tool()
def showcase_templates() -> str:
    """등록된 슬라이드 템플릿을 PNG 미리보기와 함께 쇼케이스합니다.

    각 템플릿 타입(T0~T9)의 용도, 어울리는 콘텐츠 유형, 필드 목록,
    대표 슬라이드 이미지 경로를 반환합니다.

    Returns:
        템플릿 쇼케이스 (JSON) — 이미지 경로 포함
    """
    if not DEFAULT_SLIDE_INDEX.exists():
        return "오류: slide_index.json이 없습니다. analyze_template()을 먼저 실행하세요."

    with open(DEFAULT_SLIDE_INDEX, 'r', encoding='utf-8') as f:
        idx = json.load(f)

    from collections import Counter
    tmpl_counts = Counter(s['template'] for s in idx.get('slides', []))

    SHOWCASE = {
        'T0': {
            'name': '구분페이지',
            'when': '섹션/챕터 구분, 목차 페이지',
            'content': '제목 1줄',
            'fields': ['@content_1'],
        },
        'T1': {
            'name': '카드형 다중',
            'when': '현황/문제점/개선방향, 접근전략/수행체계, 비교 분석',
            'content': '카드 2~6개 (제목 + 본문 각 1~3문장)',
            'fields': ['@governing_message', '@breadcrumb', '@content_1', '@카드N_제목', '@카드N_내용 (N=1~6)'],
        },
        'T2': {
            'name': '카드+다이어그램',
            'when': '사업 목적, 전략 개요 (카드 + 세분화 불릿)',
            'content': '카드 2~3개 + 세분화 항목',
            'fields': ['@카드N_제목/내용', '@content_N', '@label_N', '@section_title_N'],
        },
        'T3': {
            'name': '범위/개요 (거버닝메시지)',
            'when': '사업 범위, 핵심 개요, 비전',
            'content': 'Governing Message + 핵심 문구 + 6개 영역 설명',
            'fields': ['@governing_message', '@breadcrumb', '@heading_1', '@content_N', '@section_title_N'],
        },
        'T4': {
            'name': '다중 데이터테이블',
            'when': '복수 테이블, 일정표, 비교 매트릭스',
            'content': '마크다운 테이블 2개 이상',
            'fields': ['@governing_message', '마크다운 테이블'],
        },
        'T5': {
            'name': '테이블+다이어그램',
            'when': '데이터 테이블 + 시각적 설명 다이어그램',
            'content': '테이블 1개 + 설명 텍스트',
            'fields': ['@governing_message', '마크다운 테이블', '@content_N'],
        },
        'T6': {
            'name': '순수 데이터테이블',
            'when': '큰 데이터 표, 지표, 상세 스펙',
            'content': '마크다운 테이블 (10줄 이상)',
            'fields': ['@governing_message', '마크다운 테이블'],
        },
        'T7': {
            'name': '테이블+설명shape (프로세스)',
            'when': '프로세스 흐름, 단계별 설명 + 데이터',
            'content': '프로세스 단계 + 테이블',
            'fields': ['@governing_message', '@heading_N', '@content_N', '마크다운 테이블'],
        },
        'T8': {
            'name': '이미지중심',
            'when': '조직도, 시스템 구성도, 스크린샷, 사진',
            'content': '이미지 + 캡션/설명 텍스트',
            'fields': ['@governing_message', '@content_N (이미지 설명)'],
        },
        'T9': {
            'name': '핵심메시지/다이어그램',
            'when': '핵심 포인트 나열, CSF, 성공요소, 원칙',
            'content': '핵심 문구 3~6개',
            'fields': ['@governing_message', '@breadcrumb', '@content_N (N=1~6)'],
        },
    }

    result = []
    for tmpl_code in sorted(SHOWCASE.keys()):
        info = SHOWCASE[tmpl_code]
        rep_slides = [s['slide_number'] for s in idx['slides'] if s['template'] == tmpl_code][:3]
        images = []
        for sn in rep_slides:
            img = TEMPLATE_IMAGES_DIR / f"S{sn:03d}_{tmpl_code}.png"
            if img.exists():
                images.append(str(img))

        result.append({
            **info,
            'template': tmpl_code,
            'count': tmpl_counts.get(tmpl_code, 0),
            'representative_slides': rep_slides,
            'preview_images': images,
        })

    return json.dumps(result, ensure_ascii=False, indent=2)


@mcp.tool()
def match_slide(
    pptx_path: str,
    slide_number: int
) -> str:
    """PPTX 파일의 특정 슬라이드를 분석하여 기존 템플릿(T0~T9)과 매칭합니다.

    새 슬라이드를 넣으면 shape 구성(AutoShape, Table, Picture 등)을 분석하여
    가장 유사한 기존 템플릿을 찾아줍니다.

    Args:
        pptx_path: 분석할 PPTX 파일 경로
        slide_number: 분석할 슬라이드 번호 (1부터 시작)

    Returns:
        매칭 결과 (상위 5개 후보 + 특성 분석)
    """
    from template_matcher import analyze_and_match

    p = Path(pptx_path)
    if not p.exists():
        return f"오류: 파일을 찾을 수 없습니다: {pptx_path}"

    try:
        result = analyze_and_match(str(p), slide_number)

        if 'error' in result:
            return f"오류: {result['error']}"

        best = result['best_match']
        features = result['features']

        output = f"슬라이드 #{slide_number} 분석 결과\n\n"
        output += f"shape 구성: 총 {features['total_shapes']}개\n"
        output += f"  AutoShape: {features['auto_shapes']}, TextBox: {features['text_boxes']}\n"
        output += f"  카드테이블(2x1): {features['card_tables']}, 데이터테이블: {features['data_tables']}\n"
        output += f"  이미지: {features['pictures']}, 그룹: {features['groups']}\n\n"

        output += f"가장 유사한 템플릿: {best['template']} ({best['name']}) — 유사도 {best['score']:.0%}\n\n"
        output += "상위 5개 후보:\n"
        for m in result['matches']:
            output += f"  {m['template']} ({m['name']}): {m['score']:.0%}\n"

        return output

    except Exception as e:
        return f"오류: {type(e).__name__}: {e}"


@mcp.tool()
def add_template(
    pptx_path: str,
    slide_number: int,
    template_name: str = ""
) -> str:
    """새 슬라이드를 템플릿으로 등록합니다.

    지정된 PPTX의 슬라이드를 분석하여 slide_index.json에 추가하고,
    PNG 미리보기를 생성합니다.

    Args:
        pptx_path: 슬라이드가 포함된 PPTX 파일 경로
        slide_number: 등록할 슬라이드 번호
        template_name: 템플릿 이름 (예: "T11_프로세스맵"). 생략 시 자동 매칭 결과 사용

    Returns:
        등록 결과
    """
    from template_matcher import analyze_and_match
    from template_extractor import classify_shape_role, extract_shape_info
    from pptx import Presentation as PptxPresentation

    p = Path(pptx_path)
    if not p.exists():
        return f"오류: 파일을 찾을 수 없습니다: {pptx_path}"

    try:
        prs = PptxPresentation(str(p))
        if slide_number < 1 or slide_number > len(prs.slides):
            return f"오류: 슬라이드 번호 범위 초과 (1~{len(prs.slides)})"

        slide = prs.slides[slide_number - 1]

        # 템플릿 이름 자동 결정
        if not template_name:
            result = analyze_and_match(str(p), slide_number)
            best = result.get('best_match', {})
            if best.get('score', 0) >= 0.8:
                template_name = best['template']
            else:
                # 새 번호 부여
                existing = set()
                if DEFAULT_SLIDE_INDEX.exists():
                    with open(DEFAULT_SLIDE_INDEX, 'r', encoding='utf-8') as f:
                        for s in json.load(f).get('slides', []):
                            existing.add(s.get('template', ''))
                for i in range(10, 100):
                    if f'T{i}' not in existing:
                        template_name = f'T{i}'
                        break

        # shape 분석
        shapes_info = []
        role_map = {}
        for shape in slide.shapes:
            role = classify_shape_role(shape, slide.shapes)
            s_info = extract_shape_info(shape)
            s_info['role'] = role
            si = len(shapes_info)
            shapes_info.append(s_info)
            role_map.setdefault(role, []).append(si)

        # slide_index.json에 추가
        if DEFAULT_SLIDE_INDEX.exists():
            with open(DEFAULT_SLIDE_INDEX, 'r', encoding='utf-8') as f:
                idx = json.load(f)
        else:
            idx = {'slides': [], 'total_slides': 0}

        new_slide_num = idx['total_slides'] + 1
        idx['slides'].append({
            'slide_number': new_slide_num,
            'template': template_name,
            'template_name': template_name,
            'layout_name': slide.slide_layout.name,
            'shape_count': len(shapes_info),
            'role_map': role_map,
            'shapes': shapes_info,
            'source_pptx': str(p.name),
            'source_slide': slide_number,
        })
        idx['total_slides'] = new_slide_num

        with open(DEFAULT_SLIDE_INDEX, 'w', encoding='utf-8') as f:
            json.dump(idx, f, ensure_ascii=False, indent=2)

        return (
            f"템플릿 등록 완료!\n"
            f"  이름: {template_name}\n"
            f"  슬라이드 번호: {new_slide_num}\n"
            f"  shape 수: {len(shapes_info)}\n"
            f"  원본: {p.name} 슬라이드 #{slide_number}\n"
            f"  slide_index.json 업데이트됨"
        )

    except Exception as e:
        return f"오류: {type(e).__name__}: {e}"


# ---------------------------------------------------------------------------
# 진입점
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    mcp.run()
